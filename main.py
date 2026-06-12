import os
import re
import json
import glob
import time

import base64
import logging
import pptx
from concurrent.futures import ThreadPoolExecutor, as_completed

from openai import OpenAI
from jinja2 import Template
from dotenv import load_dotenv
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ==========================================
# 0. 初始化配置、日志、客户端
# ==========================================
load_dotenv()

# --- 日志配置 ---
LOG_DIR = "logs"
os.makedirs(LOG_DIR, exist_ok=True)

logger = logging.getLogger("MedAudit")
logger.setLevel(logging.DEBUG)

# 控制台 handler：INFO 级别，简洁格式
console_handler = logging.StreamHandler()
console_handler.setLevel(logging.INFO)
console_handler.setFormatter(logging.Formatter("%(message)s"))

# 文件 handler：DEBUG 级别，详细格式
file_handler = logging.FileHandler(
    os.path.join(LOG_DIR, "medaudit.log"), encoding="utf-8"
)
file_handler.setLevel(logging.DEBUG)
file_handler.setFormatter(
    logging.Formatter("%(asctime)s | %(levelname)-7s | %(message)s", datefmt="%Y-%m-%d %H:%M:%S")
)

logger.addHandler(console_handler)
logger.addHandler(file_handler)

# --- API 客户端 ---
client = OpenAI(
    base_url=os.getenv("OPENAI_BASE_URL", "https://token-plan-cn.xiaomimimo.com/v1"),
    api_key=os.getenv("OPENAI_API_KEY", ""),
)
MODEL_NAME = os.getenv("OPENAI_MODEL_ID", "mimo-v2.5-pro")
MULTIMODAL_MODEL = os.getenv("OPENAI_MULTIMODAL_MODEL", "")  # 视觉模型，为空则跳过图片识别
VISION_ENABLED = bool(MULTIMODAL_MODEL)  # 是否启用图片识别
logger.info(f"模型: {MODEL_NAME} | 视觉: {'启用 (' + MULTIMODAL_MODEL + ')' if VISION_ENABLED else '未启用'}")


# ==========================================
# 工具函数：JSON 安全解析 + LLM 重试
# ==========================================
def clean_json_string(text):
    """清理 LLM 输出中常见的 JSON 格式问题"""
    text = text.strip()
    text = re.sub(r'```json\s*', '', text)
    text = re.sub(r'```\s*$', '', text)
    text = re.sub(r'//[^\n]*', '', text)
    text = re.sub(r',\s*([}\]])', r'\1', text)
    return text.strip()


def safe_json_parse(text, retries=2, retry_messages=None, retry_prompt=None):
    """安全解析 JSON，失败时清理重试；仍失败则重调 LLM"""
    cleaned = clean_json_string(text)
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        logger.debug(f"首次 JSON 解析失败，清理后重试")

    for attempt in range(retries):
        try:
            return json.loads(cleaned)
        except json.JSONDecodeError:
            if attempt < retries - 1:
                match = re.search(r'(\[.*\]|\{.*\})', cleaned, re.DOTALL)
                if match:
                    cleaned = match.group(1)
                    logger.debug(f"提取 JSON 片段，第 {attempt + 1} 次重试解析")
                    continue
            if retry_messages and retry_prompt:
                logger.warning(f"JSON 解析失败，第 {attempt + 1} 次重调 LLM")
                retry_messages_copy = retry_messages + [
                    {"role": "user", "content": retry_prompt}
                ]
                response = client.chat.completions.create(
                    model=MODEL_NAME,
                    messages=retry_messages_copy,
                    temperature=0.1
                )
                text = response.choices[0].message.content
                cleaned = clean_json_string(text)
            else:
                raise ValueError(f"JSON 解析失败，已重试 {retries} 次。原始内容: {text[:200]}...")

    raise ValueError(f"JSON 解析失败，已重试 {retries} 次。原始内容: {text[:200]}...")


# ==========================================
# 工具函数：从规则库读取有效的药品+适应症组合
# ==========================================
def get_valid_drug_indications():
    """从 rules_config.json 读取所有有效的 药品→适应症 映射"""
    rules_file = os.path.join("config", "rules_config.json")
    if not os.path.exists(rules_file):
        logger.warning(f"规则配置文件不存在: {rules_file}")
        return {}
    with open(rules_file, "r", encoding="utf-8") as f:
        all_rules = json.load(f)

    drug_indications = {}
    for drug, indications in all_rules.items():
        drug_indications[drug] = list(indications.keys())
    return drug_indications


# ==========================================
# 1. 提取 Agent: PPT 文本 + 图片提取
# ==========================================
def extract_text_from_pptx(pptx_path):
    """提取PPT内所有文本框和表格数据"""
    logger.debug(f"开始提取PPT文本: {pptx_path}")
    prs = pptx.Presentation(pptx_path)
    extracted_texts = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_text = f"--- Slide {slide_idx + 1} --- \n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_text += paragraph.text.strip() + "\n"
            if shape.has_table:
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_data:
                        slide_text += " [表格数据]: " + " | ".join(row_data) + "\n"
        extracted_texts.append(slide_text)

    logger.debug(f"提取完成，共 {len(prs.slides)} 页幻灯片")
    return "\n".join(extracted_texts)


def extract_images_from_pptx(pptx_path, output_dir):
    """提取PPT中所有嵌入的图片，保存到指定目录"""
    logger.debug(f"开始提取PPT图片: {pptx_path}")
    prs = pptx.Presentation(pptx_path)
    filename = os.path.splitext(os.path.basename(pptx_path))[0]
    img_dir = os.path.join(output_dir, "images", filename)
    os.makedirs(img_dir, exist_ok=True)

    image_info_list = []
    skipped_count = 0

    for slide_idx, slide in enumerate(prs.slides):
        img_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_count += 1
                    image = shape.image
                    ext = image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    # 跳过不支持的格式
                    if ext.lower() not in ['bmp', 'gif', 'jpg', 'png', 'tiff', 'wmf']:
                        skipped_count += 1
                        logger.debug(f"  跳过不支持的图片格式: {ext}")
                        continue
                    img_filename = f"slide_{slide_idx + 1}_img_{img_count}.{ext}"
                    img_path = os.path.join(img_dir, img_filename)

                    with open(img_path, "wb") as f:
                        f.write(image.blob)

                    image_info_list.append((img_path, slide_idx + 1))
                    logger.debug(f"  保存图片: {img_filename}")
                except ValueError as e:
                    skipped_count += 1
                    logger.debug(f"  跳过图片: {e}")
                    continue

    # 返回图片列表和日志行（由调用方统一输出）
    if skipped_count > 0:
        log_line = f"  │  提取图片 {len(image_info_list)} 张（跳过 {skipped_count} 张不支持的格式）"
    else:
        log_line = f"  │  提取图片 {len(image_info_list)} 张"
    return image_info_list, log_line


def image_to_base64(image_path):
    """将图片文件转为 base64 编码"""
    with open(image_path, "rb") as f:
        return base64.b64encode(f.read()).decode("utf-8")


def describe_images_with_llm(image_info_list):
    """利用多模态 LLM 识别图片内容，返回描述文本"""
    if not image_info_list or not VISION_ENABLED:
        if image_info_list and not VISION_ENABLED:
            logger.warning("    视觉模型未配置，跳过图片识别。设置 OPENAI_MULTIMODAL_MODEL 环境变量以启用。")
        return ""

    descriptions = []
    for img_path, slide_idx in image_info_list:
        try:
            img_b64 = image_to_base64(img_path)
            ext = os.path.splitext(img_path)[1].lower()
            mime_map = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".gif": "image/gif", ".webp": "image/webp"}
            mime_type = mime_map.get(ext, "image/png")

            logger.debug(f"  调用视觉模型识别: Slide {slide_idx} 图片")
            response = client.chat.completions.create(
                model=MULTIMODAL_MODEL,
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "请仔细识别这张医学病例PPT中的图片内容。如果是化验单/检查报告，请提取关键数值（如EOS、FEV1等）；如果是CT/MRI影像，请描述所见；如果包含患者个人信息（姓名、住院号、医生信息），请特别标注。用中文简洁描述。"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{img_b64}"
                            }
                        }
                    ]
                }],
                temperature=0.1
            )
            desc = response.choices[0].message.content.strip()
            descriptions.append(f"[Slide {slide_idx} 图片内容]: {desc}")
            logger.debug(f"  Slide {slide_idx} 图片识别完成: {desc[:50]}...")
        except Exception as e:
            logger.error(f"  Slide {slide_idx} 图片识别失败: {e}")
            descriptions.append(f"[Slide {slide_idx} 图片识别失败]: {str(e)}")

    return "\n".join(descriptions)


def transform_raw_text_to_structured(raw_text, image_desc_text=""):
    """利用LLM将PPT文本+图片描述转化为结构化的患者信息JSON"""
    combined_text = raw_text
    if image_desc_text:
        combined_text += "\n\n【图片识别内容】:\n" + image_desc_text

    drug_indications = get_valid_drug_indications()
    constraint_str = ""
    if drug_indications:
        drug_list = list(drug_indications.keys())
        # 构建候选适应症说明（过滤掉 default/unknown 这种非医学适应症）
        indication_options = {}
        for drug, inds in drug_indications.items():
            real_inds = [i for i in inds if i not in ("default", "unknown")]
            indication_options[drug] = real_inds if real_inds else inds
        constraint_str = f"""
    【重要约束】：
    - "treatment_drug" 只能填药品通用名，必须从以下列表中选择：{json.dumps(drug_list, ensure_ascii=False)}。如不在列表中，填实际药品名。
    - "indication" 填该病例实际对应的医学适应症，从该药品的候选适应症中选择：
      {json.dumps(indication_options, ensure_ascii=False, indent=6)}
      注意：default 不是医学适应症，不要使用。如确实无法匹配，留空字符串 ""。
    - 注意：treatment_drug 只填药品名（如"西甲硅油"），不要包含适应症信息。
    """

    prompt = f"""
    你是一个精通临床医学数据处理的AI。请阅读以下从病例幻灯片(PPT)中提取出来的原始文本及图片识别内容，将其清洗并提取为标准的结构化数据。
    {constraint_str}
    【原始文本】:
    {combined_text}

    【输出格式要求】:
    必须只输出一个合法的 JSON 对象，不要包含任何 Markdown 标记（如 ```json）。字段要求如下：
    {{
        "patient_id": "患者编号或匿名化代号",
        "gender": "男/女/不详",
        "age": "年龄数字",
        "admission_date": "入院日期/就诊日期(YYYY-MM-DD)",
        "chief_complaint": "主诉",
        "present_illness": "现病史",
        "past_history": "既往史/个人史",
        "diagnosis": "初步诊断/最终诊断",
        "treatment_drug": "药品通用名（仅填药品名，如：西甲硅油、美泊利珠单抗、氟替美维）",
        "treatment_detail": "具体用法用量及治疗过程描述",
        "follow_up": "随访结果、检查指标及转归描述",
        "indication": "该病例实际对应的医学适应症（如：功能性消化不良、肠易激综合征、SEA、COPD、CRSwNP、EGPA、哮喘等，无法匹配则留空）",
        "privacy_issues": "如果图片或文本中发现患者个人信息泄露（姓名、住院号、医生信息等），在此列出；没有则填'无'"
    }}
    """
    logger.debug("调用 LLM 进行结构化数据提取")
    messages = [{"role": "user", "content": prompt}]
    response = client.chat.completions.create(
        model=MODEL_NAME,
        messages=messages,
        temperature=0.1
    )
    content = response.choices[0].message.content
    logger.debug(f"LLM 结构化原始输出长度: {len(content)} 字符")
    result = safe_json_parse(
        content,
        retries=2,
        retry_messages=messages,
        retry_prompt="你上次输出的不是合法JSON，请严格只输出JSON对象，不要有任何其他文字。"
    )

    # 后处理：清理 treatment_drug 字段，去掉混入的适应症信息
    if "treatment_drug" in result:
        drug = result["treatment_drug"]
        # 去掉 " - xxx" 后缀（如 "西甲硅油 - default" → "西甲硅油"）
        if " - " in drug:
            result["treatment_drug"] = drug.split(" - ")[0].strip()
        # 去掉括号内容（如 "美泊利珠单抗(SEA)" → "美泊利珠单抗"）
        if "(" in result["treatment_drug"]:
            result["treatment_drug"] = result["treatment_drug"].split("(")[0].strip()

    # 后处理：indication 不应是药品名
    if "indication" in result:
        ind = result["indication"]
        # 如果 indication 和 treatment_drug 相同，说明 LLM 混淆了
        if ind == result.get("treatment_drug"):
            result["indication"] = "unknown"

    return result


# ==========================================
# 2. 审核 Agent: 匹配规则与合规审查
# ==========================================
def load_rules(drug_name, indication=None):
    """从本地 JSON 配置中加载特定药品+适应症的审核规则"""
    rules_file = os.path.join("config", "rules_config.json")
    if not os.path.exists(rules_file):
        logger.warning(f"规则配置文件不存在: {rules_file}")
        return None
    with open(rules_file, "r", encoding="utf-8") as f:
        all_rules = json.load(f)

    matched_drug = None
    for key in all_rules.keys():
        if key in drug_name or drug_name in key:
            matched_drug = all_rules[key]
            logger.debug(f"规则匹配: 药品 '{drug_name}' → '{key}'")
            break

    if not matched_drug:
        logger.debug(f"未找到药品 '{drug_name}' 的审核规则")
        return None

    if indication:
        for ind_key in matched_drug.keys():
            if ind_key in indication or indication in ind_key:
                logger.debug(f"规则匹配: 适应症 '{indication}' → '{ind_key}'")
                return matched_drug[ind_key]

    # 未匹配到适应症，返回 None（不兜底）
    logger.debug(f"未匹配到适应症 '{indication}' 的专属规则，使用通用审核")
    return None


def audit_patient_data(patient_data, rules):
    """对照特殊规则，让大模型输出结构化审核意见"""
    if not rules:
        rules_str = "无特殊药品专属规则，仅按照通用规范审核：1. 保护患者隐私；2. 逻辑故事线通顺；3. 药物需用通用名。"
        logger.info("    无专属规则，使用通用规范审核")
    else:
        rules_str = json.dumps(rules, ensure_ascii=False, indent=2)
        logger.debug(f"    加载审核规则: {len(rules_str)} 字符")

    prompt = f"""
    你是一名资深的药企医学合规经理。请严格对照【审核规则】，对【患者结构化数据】进行逐条审查，指出不合规或存在安全隐患的地方。

    【患者结构化数据】:
    {json.dumps(patient_data, ensure_ascii=False, indent=2)}

    【审核规则】:
    {rules_str}

    【输出格式要求】:
    请仅返回一个标准的 JSON 数组，不要包裹 ```json 标记。数组中每个对象代表一个审核项。

    **重要：必须包含以下【基础质控检查】的4个固定审核维度，放在数组最前面：**

    1. **文本规范性检查** - 检查病例文本中是否存在错别字、语法错误、标点符号误用等问题。如发现错别字（如"的地得"混用、医学术语拼写错误等），给出具体位置和修正建议。
    2. **通用名规范检查** - 检查药品名称是否使用了通用名（INN），而非商品名或品牌名。如发现使用商品名（如"舒利迭"应为"沙美特罗替卡松"、"信必可"应为"布地奈德福莫特罗"等），需标注并建议替换为通用名。
    3. **医学逻辑审核** - 检查药物与适应症的医学逻辑是否合理，包括：药物作用机制与疾病病理是否匹配、剂量是否在合理范围、用药疗程是否符合临床实践、联合用药是否存在相互作用风险等。
    4. **单位规范检查** - 检查医学数值的单位是否正确规范，包括：剂量单位（mg/mg/mL）、给药频率单位（qd/bid/tid/q4w等）、检验指标单位（%、×10^9/L等）、体重/身高等单位。如发现单位缺失、错误或不规范，需标注并给出正确写法。

    然后再根据【审核规则】进行其他维度的审核。

    格式如下：
    [
        {{
            "aspect": "文本规范性检查",
            "status": "Pass 或 Warning 或 Fail",
            "detail": "具体的审查判定说明",
            "suggestion": "修改建议或'无需修改'"
        }},
        {{
            "aspect": "通用名规范检查",
            "status": "Pass 或 Warning 或 Fail",
            "detail": "具体的审查判定说明",
            "suggestion": "修改建议或'无需修改'"
        }},
        {{
            "aspect": "医学逻辑审核",
            "status": "Pass 或 Warning 或 Fail",
            "detail": "具体的审查判定说明",
            "suggestion": "修改建议或'无需修改'"
        }},
        {{
            "aspect": "单位规范检查",
            "status": "Pass 或 Warning 或 Fail",
            "detail": "具体的审查判定说明",
            "suggestion": "修改建议或'无需修改'"
        }},
        {{
            "aspect": "其他审核维度(如: 适应症合规/隐私保护/Storyline故事线/安全性风险/科学幻灯)",
            "status": "Pass 或 Warning 或 Fail",
            "detail": "具体的审查判定说明，引用病例中的具体事实",
            "suggestion": "修改建议或'无需修改'"
        }}
    ]
    """
    logger.debug("调用 LLM 进行合规审核")
    messages = [{"role": "user", "content": prompt}]
    response = client.chat.completions.create(
        model=MODEL_NAME,
        messages=messages,
        temperature=0.2
    )
    content = response.choices[0].message.content
    logger.debug(f"LLM 审核原始输出长度: {len(content)} 字符")
    return safe_json_parse(
        content,
        retries=2,
        retry_messages=messages,
        retry_prompt="你上次输出的不是合法JSON数组，请严格只输出JSON数组，不要有任何其他文字。"
    )


# ==========================================
# 3. 报告生成 Agent: HTML 模板渲染
# ==========================================
HTML_TEMPLATE = """
<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>病例医学合规审核报告</title>
    <link href="https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@300;400;500;600;700&family=JetBrains+Mono:wght@400;500&display=swap" rel="stylesheet">
    <style>
        :root {
            --color-bg: #f8f9fc;
            --color-surface: #ffffff;
            --color-border: #e2e8f0;
            --color-text-primary: #1a202c;
            --color-text-secondary: #4a5568;
            --color-text-muted: #718096;
            --color-accent: #2b6cb0;
            --color-accent-light: #ebf4ff;
            --color-pass: #059669;
            --color-pass-bg: #ecfdf5;
            --color-pass-border: #a7f3d0;
            --color-warning: #d97706;
            --color-warning-bg: #fffbeb;
            --color-warning-border: #fde68a;
            --color-fail: #dc2626;
            --color-fail-bg: #fef2f2;
            --color-fail-border: #fecaca;
            --color-privacy: #dc2626;
            --color-privacy-bg: #fef2f2;
            --radius-sm: 6px;
            --radius-md: 10px;
            --radius-lg: 16px;
            --shadow-sm: 0 1px 3px rgba(0,0,0,0.06), 0 1px 2px rgba(0,0,0,0.04);
            --shadow-md: 0 4px 12px rgba(0,0,0,0.08), 0 2px 4px rgba(0,0,0,0.04);
            --shadow-lg: 0 12px 40px rgba(0,0,0,0.1), 0 4px 12px rgba(0,0,0,0.05);
        }

        * { margin: 0; padding: 0; box-sizing: border-box; }

        body {
            font-family: 'Noto Sans SC', -apple-system, BlinkMacSystemFont, sans-serif;
            background: var(--color-bg);
            color: var(--color-text-primary);
            line-height: 1.7;
            min-height: 100vh;
            padding: 40px 20px;
            background-image:
                radial-gradient(ellipse at 20% 0%, rgba(43, 108, 176, 0.04) 0%, transparent 60%),
                radial-gradient(ellipse at 80% 100%, rgba(5, 150, 105, 0.03) 0%, transparent 60%);
        }

        .report-wrapper {
            max-width: 960px;
            margin: 0 auto;
        }

        /* ===== HEADER ===== */
        .report-header {
            background: linear-gradient(135deg, #1e3a5f 0%, #2b6cb0 100%);
            color: white;
            padding: 48px 48px 40px;
            border-radius: var(--radius-lg) var(--radius-lg) 0 0;
            position: relative;
            overflow: hidden;
        }
        .report-header::before {
            content: '';
            position: absolute;
            top: -50%;
            right: -20%;
            width: 400px;
            height: 400px;
            background: radial-gradient(circle, rgba(255,255,255,0.08) 0%, transparent 70%);
            pointer-events: none;
        }
        .report-header::after {
            content: '';
            position: absolute;
            bottom: 0;
            left: 0;
            right: 0;
            height: 4px;
            background: linear-gradient(90deg, #059669, #2b6cb0, #d97706);
        }
        .header-label {
            font-size: 11px;
            font-weight: 600;
            letter-spacing: 3px;
            text-transform: uppercase;
            opacity: 0.7;
            margin-bottom: 12px;
        }
        .header-title {
            font-size: 28px;
            font-weight: 700;
            letter-spacing: -0.5px;
            margin-bottom: 8px;
        }
        .header-subtitle {
            font-size: 14px;
            opacity: 0.65;
            font-weight: 300;
        }

        /* ===== MAIN BODY ===== */
        .report-body {
            background: var(--color-surface);
            padding: 0 48px 48px;
            border-radius: 0 0 var(--radius-lg) var(--radius-lg);
            box-shadow: var(--shadow-lg);
        }

        /* ===== PATIENT INFO ===== */
        .patient-section {
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 24px;
            padding: 36px 0;
            border-bottom: 1px solid var(--color-border);
        }
        .info-block {
            display: flex;
            flex-direction: column;
            gap: 4px;
        }
        .info-label {
            font-size: 11px;
            font-weight: 600;
            letter-spacing: 1.5px;
            text-transform: uppercase;
            color: var(--color-text-muted);
        }
        .info-value {
            font-size: 15px;
            font-weight: 500;
            color: var(--color-text-primary);
        }
        .info-value.drug {
            color: var(--color-accent);
            font-weight: 600;
        }
        .info-value.indication {
            display: inline-flex;
            align-items: center;
            gap: 6px;
            background: var(--color-accent-light);
            color: var(--color-accent);
            padding: 4px 12px;
            border-radius: 20px;
            font-size: 13px;
            font-weight: 600;
            width: fit-content;
        }

        /* ===== PRIVACY ALERT ===== */
        .privacy-alert {
            margin-top: 24px;
            padding: 16px 20px;
            background: var(--color-privacy-bg);
            border: 1px solid var(--color-fail-border);
            border-radius: var(--radius-md);
            display: flex;
            align-items: flex-start;
            gap: 12px;
        }
        .privacy-alert .alert-icon {
            font-size: 18px;
            flex-shrink: 0;
            margin-top: 1px;
        }
        .privacy-alert .alert-content {
            font-size: 13px;
            color: var(--color-privacy);
            line-height: 1.6;
        }
        .privacy-alert .alert-title {
            font-weight: 600;
            margin-bottom: 4px;
        }

        /* ===== SUMMARY STATS ===== */
        .summary-bar {
            display: flex;
            gap: 12px;
            padding: 28px 0;
            border-bottom: 1px solid var(--color-border);
        }
        .stat-chip {
            display: flex;
            align-items: center;
            gap: 8px;
            padding: 8px 16px;
            border-radius: 24px;
            font-size: 13px;
            font-weight: 600;
        }
        .stat-chip.pass {
            background: var(--color-pass-bg);
            color: var(--color-pass);
            border: 1px solid var(--color-pass-border);
        }
        .stat-chip.warning {
            background: var(--color-warning-bg);
            color: var(--color-warning);
            border: 1px solid var(--color-warning-border);
        }
        .stat-chip.fail {
            background: var(--color-fail-bg);
            color: var(--color-fail);
            border: 1px solid var(--color-fail-border);
        }
        .stat-dot {
            width: 8px;
            height: 8px;
            border-radius: 50%;
        }
        .stat-chip.pass .stat-dot { background: var(--color-pass); }
        .stat-chip.warning .stat-dot { background: var(--color-warning); }
        .stat-chip.fail .stat-dot { background: var(--color-fail); }

        /* ===== AUDIT TABLE ===== */
        .section-title {
            font-size: 16px;
            font-weight: 600;
            color: var(--color-text-primary);
            padding: 32px 0 20px;
            display: flex;
            align-items: center;
            gap: 10px;
        }
        .section-title::before {
            content: '';
            width: 4px;
            height: 20px;
            background: var(--color-accent);
            border-radius: 2px;
        }
        .subsection-title {
            font-size: 14px;
            font-weight: 600;
            color: var(--color-text-secondary);
            padding: 20px 0 12px;
            display: flex;
            align-items: center;
            gap: 8px;
        }
        .qc-table {
            border: 2px solid #dbeafe;
            background: #f0f7ff;
        }
        .qc-table thead th {
            background: #dbeafe;
            color: #1e40af;
        }

        .audit-table {
            width: 100%;
            border-collapse: separate;
            border-spacing: 0;
            border: 1px solid var(--color-border);
            border-radius: var(--radius-md);
            overflow: hidden;
        }
        .audit-table thead th {
            background: #f7fafc;
            color: var(--color-text-secondary);
            font-size: 11px;
            font-weight: 600;
            letter-spacing: 1px;
            text-transform: uppercase;
            padding: 14px 20px;
            text-align: left;
            border-bottom: 2px solid var(--color-border);
        }
        .audit-table tbody tr {
            transition: background 0.15s ease;
        }
        .audit-table tbody tr:hover {
            background: #f7fafc;
        }
        .audit-table tbody tr:not(:last-child) td {
            border-bottom: 1px solid var(--color-border);
        }
        .audit-table td {
            padding: 14px 16px;
            font-size: 13px;
            vertical-align: top;
            word-break: break-word;
            line-height: 1.7;
        }
        .audit-table td:first-child {
            font-weight: 600;
            color: var(--color-text-primary);
            white-space: nowrap;
            width: 10%;
        }
        .audit-table td:nth-child(2) {
            width: 8%;
            text-align: center;
        }
        .audit-table td:nth-child(3) {
            color: var(--color-text-secondary);
            width: 42%;
        }
        .audit-table td:nth-child(4) {
            color: var(--color-text-secondary);
            width: 40%;
        }

        /* Status badges */
        .status-badge {
            display: inline-flex;
            align-items: center;
            gap: 6px;
            padding: 4px 12px;
            border-radius: 20px;
            font-size: 12px;
            font-weight: 600;
            letter-spacing: 0.5px;
            white-space: nowrap;
        }
        .status-badge::before {
            content: '';
            width: 6px;
            height: 6px;
            border-radius: 50%;
        }
        .status-Pass {
            background: var(--color-pass-bg);
            color: var(--color-pass);
            border: 1px solid var(--color-pass-border);
        }
        .status-Pass::before { background: var(--color-pass); }
        .status-Warning {
            background: var(--color-warning-bg);
            color: var(--color-warning);
            border: 1px solid var(--color-warning-border);
        }
        .status-Warning::before { background: var(--color-warning); }
        .status-Fail {
            background: var(--color-fail-bg);
            color: var(--color-fail);
            border: 1px solid var(--color-fail-border);
        }
        .status-Fail::before { background: var(--color-fail); }

        /* ===== FOOTER ===== */
        .report-footer {
            margin-top: 32px;
            padding-top: 24px;
            border-top: 1px solid var(--color-border);
            display: flex;
            justify-content: space-between;
            align-items: center;
            font-size: 11px;
            color: var(--color-text-muted);
        }
        .footer-logo {
            font-weight: 600;
            letter-spacing: 1px;
        }

        /* ===== RESPONSIVE ===== */
        @media (max-width: 768px) {
            body { padding: 16px; }
            .report-header { padding: 32px 24px; }
            .report-body { padding: 0 24px 32px; }
            .header-title { font-size: 22px; }
            .patient-section { grid-template-columns: 1fr; gap: 16px; }
            .summary-bar { flex-wrap: wrap; }
            .audit-table td, .audit-table th { padding: 12px 14px; }
            .audit-table td:first-child { white-space: normal; }
            .audit-table td:nth-child(2) { width: auto; }
        }

        /* ===== PRINT ===== */
        @media print {
            body { background: white; padding: 0; }
            .report-wrapper { max-width: 100%; }
            .report-header { background: #1e3a5f !important; -webkit-print-color-adjust: exact; print-color-adjust: exact; }
            .report-body { box-shadow: none; }
            .audit-table tbody tr:hover { background: transparent; }
            .no-print { display: none !important; }
        }
    </style>
</head>
<body>
<div class="no-print" style="position: fixed; top: 20px; right: 20px; z-index: 1000;">
    <button onclick="window.print()" style="
        background: linear-gradient(135deg, #1e3a5f 0%, #2b6cb0 100%);
        color: white;
        border: none;
        padding: 12px 24px;
        border-radius: 8px;
        font-size: 14px;
        font-weight: 600;
        cursor: pointer;
        box-shadow: 0 4px 12px rgba(43, 108, 176, 0.3);
        display: flex;
        align-items: center;
        gap: 8px;
        transition: all 0.2s ease;
    " onmouseover="this.style.transform='translateY(-2px)';this.style.boxShadow='0 6px 16px rgba(43, 108, 176,0.4)'" onmouseout="this.style.transform='translateY(0)';this.style.boxShadow='0 4px 12px rgba(43, 108, 176,0.3)'">
        <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
            <path d="M6 9V2h12v7"></path>
            <path d="M6 18H4a2 2 0 0 1-2-2v-5a2 2 0 0 1 2-2h16a2 2 0 0 1 2 2v5a2 2 0 0 1-2 2h-2"></path>
            <rect x="6" y="14" width="12" height="8"></rect>
        </svg>
        导出PDF
    </button>
</div>
<div class="report-wrapper">

    <!-- Header -->
    <header class="report-header">
        <div class="header-label">Medical Compliance Audit</div>
        <h1 class="header-title">病例医学合规自动审核意见书</h1>
        <p class="header-subtitle">由 MedAudit Agent 自动生成 · 基于药品规则库的多维度合规审查</p>
    </header>

    <!-- Body -->
    <main class="report-body">

        <!-- Patient Info -->
        <section class="patient-section">
            <div class="info-block">
                <span class="info-label">患者编号</span>
                <span class="info-value">{{ patient.patient_id }}</span>
            </div>
            <div class="info-block">
                <span class="info-label">性别 / 年龄</span>
                <span class="info-value">{{ patient.gender }} / {{ patient.age }} 岁</span>
            </div>
            <div class="info-block">
                <span class="info-label">主要诊断</span>
                <span class="info-value">{{ patient.diagnosis }}</span>
            </div>
            <div class="info-block">
                <span class="info-label">治疗药物</span>
                <span class="info-value drug">{{ patient.treatment_drug }}</span>
            </div>
            <div class="info-block">
                <span class="info-label">适应症</span>
                <span class="info-value indication">{{ patient.indication }}</span>
            </div>
        </section>

        {% if patient.privacy_issues and patient.privacy_issues != '无' %}
        <div class="privacy-alert">
            <span class="alert-icon">🔒</span>
            <div class="alert-content">
                <div class="alert-title">隐私泄露风险检测</div>
                <div>{{ patient.privacy_issues }}</div>
            </div>
        </div>
        {% endif %}

        <!-- Summary Stats -->
        <div class="summary-bar">
            {% set pass_count = audit_results|selectattr('status', 'equalto', 'Pass')|list|length %}
            {% set warn_count = audit_results|selectattr('status', 'equalto', 'Warning')|list|length %}
            {% set fail_count = audit_results|selectattr('status', 'equalto', 'Fail')|list|length %}
            <div class="stat-chip pass"><span class="stat-dot"></span>Pass {{ pass_count }}</div>
            <div class="stat-chip warning"><span class="stat-dot"></span>Warning {{ warn_count }}</div>
            <div class="stat-chip fail"><span class="stat-dot"></span>Fail {{ fail_count }}</div>
        </div>

        <!-- Audit Results -->
        <h2 class="section-title">详细审核意见</h2>

        {% set qc_aspects = ['文本规范性检查', '通用名规范检查', '医学逻辑审核', '单位规范检查'] %}
        {% set qc_results = audit_results|selectattr('aspect', 'in', qc_aspects)|list %}
        {% set other_results = audit_results|rejectattr('aspect', 'in', qc_aspects)|list %}

        {% if qc_results %}
        <h3 class="subsection-title">🔍 基础质控检查</h3>
        <table class="audit-table qc-table">
            <thead>
                <tr>
                    <th style="width: 10%;">检查项目</th>
                    <th style="width: 8%;">状态</th>
                    <th style="width: 42%;">检查详情</th>
                    <th style="width: 40%;">修改建议</th>
                </tr>
            </thead>
            <tbody>
                {% for item in qc_results %}
                <tr>
                    <td>{{ item.aspect }}</td>
                    <td><span class="status-badge status-{{ item.status }}">{{ item.status }}</span></td>
                    <td>{{ item.detail }}</td>
                    <td>{{ item.suggestion }}</td>
                </tr>
                {% endfor %}
            </tbody>
        </table>
        {% endif %}

        {% if other_results %}
        <h3 class="subsection-title">📋 合规专项审核</h3>
        <table class="audit-table">
            <thead>
                <tr>
                    <th style="width: 10%;">审核维度</th>
                    <th style="width: 8%;">状态</th>
                    <th style="width: 42%;">依据详情</th>
                    <th style="width: 40%;">优化建议</th>
                </tr>
            </thead>
            <tbody>
                {% for item in other_results %}
                <tr>
                    <td>{{ item.aspect }}</td>
                    <td><span class="status-badge status-{{ item.status }}">{{ item.status }}</span></td>
                    <td>{{ item.detail }}</td>
                    <td>{{ item.suggestion }}</td>
                </tr>
                {% endfor %}
            </tbody>
        </table>
        {% endif %}

        <!-- Footer -->
        <footer class="report-footer">
            <span class="footer-logo">MedAudit Agent</span>
            <span>本报告由 AI 自动生成，仅供合规审核参考</span>
        </footer>

    </main>
</div>
</body>
</html>
"""


def generate_report(patient_data, audit_results, output_path):
    """利用 Jinja2 渲染精美的红绿灯状态审核报告"""
    template = Template(HTML_TEMPLATE)
    html_out = template.render(patient=patient_data, audit_results=audit_results)

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html_out)
    logger.info(f"  审核报告已生成: {output_path}")


COMBINED_HTML_TEMPLATE = """
<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>病例医学合规批量审核报告</title>
    <link href="https://fonts.googleapis.com/css2?family=Noto+Sans+SC:wght@300;400;500;600;700&family=JetBrains+Mono:wght@400;500&display=swap" rel="stylesheet">
    <style>
        :root {
            --color-bg: #f8f9fc;
            --color-surface: #ffffff;
            --color-border: #e2e8f0;
            --color-text-primary: #1a202c;
            --color-text-secondary: #4a5568;
            --color-text-muted: #718096;
            --color-accent: #2b6cb0;
            --color-accent-light: #ebf4ff;
            --color-pass: #059669;
            --color-pass-bg: #ecfdf5;
            --color-pass-border: #a7f3d0;
            --color-warning: #d97706;
            --color-warning-bg: #fffbeb;
            --color-warning-border: #fde68a;
            --color-fail: #dc2626;
            --color-fail-bg: #fef2f2;
            --color-fail-border: #fecaca;
            --color-privacy: #dc2626;
            --color-privacy-bg: #fef2f2;
            --radius-sm: 6px;
            --radius-md: 10px;
            --radius-lg: 16px;
            --shadow-sm: 0 1px 3px rgba(0,0,0,0.06), 0 1px 2px rgba(0,0,0,0.04);
            --shadow-md: 0 4px 12px rgba(0,0,0,0.08), 0 2px 4px rgba(0,0,0,0.04);
            --shadow-lg: 0 12px 40px rgba(0,0,0,0.1), 0 4px 12px rgba(0,0,0,0.05);
        }

        * { margin: 0; padding: 0; box-sizing: border-box; }

        body {
            font-family: 'Noto Sans SC', -apple-system, BlinkMacSystemFont, sans-serif;
            background: var(--color-bg);
            color: var(--color-text-primary);
            line-height: 1.7;
            min-height: 100vh;
            padding: 40px 20px;
            background-image:
                radial-gradient(ellipse at 20% 0%, rgba(43, 108, 176, 0.04) 0%, transparent 60%),
                radial-gradient(ellipse at 80% 100%, rgba(5, 150, 105, 0.03) 0%, transparent 60%);
        }

        .report-wrapper {
            max-width: 960px;
            margin: 0 auto;
        }

        /* ===== HEADER ===== */
        .report-header {
            background: linear-gradient(135deg, #1e3a5f 0%, #2b6cb0 100%);
            color: white;
            padding: 48px 48px 40px;
            border-radius: var(--radius-lg) var(--radius-lg) 0 0;
            position: relative;
            overflow: hidden;
        }
        .report-header::before {
            content: '';
            position: absolute;
            top: -50%;
            right: -20%;
            width: 400px;
            height: 400px;
            background: radial-gradient(circle, rgba(255,255,255,0.08) 0%, transparent 70%);
            pointer-events: none;
        }
        .report-header::after {
            content: '';
            position: absolute;
            bottom: 0;
            left: 0;
            right: 0;
            height: 4px;
            background: linear-gradient(90deg, #059669, #2b6cb0, #d97706);
        }
        .header-label {
            font-size: 11px;
            font-weight: 600;
            letter-spacing: 3px;
            text-transform: uppercase;
            opacity: 0.7;
            margin-bottom: 12px;
        }
        .header-title {
            font-size: 28px;
            font-weight: 700;
            letter-spacing: -0.5px;
            margin-bottom: 8px;
        }
        .header-subtitle {
            font-size: 14px;
            opacity: 0.65;
            font-weight: 300;
        }

        /* ===== MAIN BODY ===== */
        .report-body {
            background: var(--color-surface);
            padding: 0 48px 48px;
            border-radius: 0 0 var(--radius-lg) var(--radius-lg);
            box-shadow: var(--shadow-lg);
        }

        /* ===== GLOBAL SUMMARY ===== */
        .global-summary {
            display: flex;
            gap: 12px;
            padding: 28px 0;
            border-bottom: 1px solid var(--color-border);
            flex-wrap: wrap;
        }
        .stat-chip {
            display: flex;
            align-items: center;
            gap: 8px;
            padding: 8px 16px;
            border-radius: 24px;
            font-size: 13px;
            font-weight: 600;
        }
        .stat-chip.pass {
            background: var(--color-pass-bg);
            color: var(--color-pass);
            border: 1px solid var(--color-pass-border);
        }
        .stat-chip.warning {
            background: var(--color-warning-bg);
            color: var(--color-warning);
            border: 1px solid var(--color-warning-border);
        }
        .stat-chip.fail {
            background: var(--color-fail-bg);
            color: var(--color-fail);
            border: 1px solid var(--color-fail-border);
        }
        .stat-chip.total {
            background: var(--color-accent-light);
            color: var(--color-accent);
            border: 1px solid #bee3f8;
        }
        .stat-dot {
            width: 8px;
            height: 8px;
            border-radius: 50%;
        }
        .stat-chip.pass .stat-dot { background: var(--color-pass); }
        .stat-chip.warning .stat-dot { background: var(--color-warning); }
        .stat-chip.fail .stat-dot { background: var(--color-fail); }
        .stat-chip.total .stat-dot { background: var(--color-accent); }

        /* ===== CASE DIVIDER ===== */
        .case-section {
            padding-top: 36px;
            margin-top: 36px;
            border-top: 2px solid var(--color-border);
        }
        .case-section:first-of-type {
            margin-top: 0;
            padding-top: 28px;
            border-top: none;
        }
        .case-header {
            display: flex;
            align-items: center;
            gap: 12px;
            margin-bottom: 24px;
        }
        .case-number {
            display: inline-flex;
            align-items: center;
            justify-content: center;
            width: 32px;
            height: 32px;
            background: var(--color-accent);
            color: white;
            border-radius: 50%;
            font-size: 14px;
            font-weight: 700;
            flex-shrink: 0;
        }
        .case-title {
            font-size: 18px;
            font-weight: 600;
            color: var(--color-text-primary);
        }

        /* ===== PATIENT INFO ===== */
        .patient-section {
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 24px;
            padding-bottom: 24px;
            border-bottom: 1px solid var(--color-border);
        }
        .info-block {
            display: flex;
            flex-direction: column;
            gap: 4px;
        }
        .info-label {
            font-size: 11px;
            font-weight: 600;
            letter-spacing: 1.5px;
            text-transform: uppercase;
            color: var(--color-text-muted);
        }
        .info-value {
            font-size: 15px;
            font-weight: 500;
            color: var(--color-text-primary);
        }
        .info-value.drug {
            color: var(--color-accent);
            font-weight: 600;
        }
        .info-value.indication {
            display: inline-flex;
            align-items: center;
            gap: 6px;
            background: var(--color-accent-light);
            color: var(--color-accent);
            padding: 4px 12px;
            border-radius: 20px;
            font-size: 13px;
            font-weight: 600;
            width: fit-content;
        }

        /* ===== PRIVACY ALERT ===== */
        .privacy-alert {
            margin-top: 24px;
            padding: 16px 20px;
            background: var(--color-privacy-bg);
            border: 1px solid var(--color-fail-border);
            border-radius: var(--radius-md);
            display: flex;
            align-items: flex-start;
            gap: 12px;
        }
        .privacy-alert .alert-icon {
            font-size: 18px;
            flex-shrink: 0;
            margin-top: 1px;
        }
        .privacy-alert .alert-content {
            font-size: 13px;
            color: var(--color-privacy);
            line-height: 1.6;
        }
        .privacy-alert .alert-title {
            font-weight: 600;
            margin-bottom: 4px;
        }

        /* ===== CASE SUMMARY ===== */
        .case-summary {
            display: flex;
            gap: 12px;
            padding: 20px 0;
            flex-wrap: wrap;
        }

        /* ===== AUDIT TABLE ===== */
        .section-title {
            font-size: 16px;
            font-weight: 600;
            color: var(--color-text-primary);
            padding: 20px 0 16px;
            display: flex;
            align-items: center;
            gap: 10px;
        }
        .section-title::before {
            content: '';
            width: 4px;
            height: 20px;
            background: var(--color-accent);
            border-radius: 2px;
        }
        .subsection-title {
            font-size: 14px;
            font-weight: 600;
            color: var(--color-text-secondary);
            padding: 16px 0 10px;
            display: flex;
            align-items: center;
            gap: 8px;
        }
        .qc-table {
            border: 2px solid #dbeafe;
            background: #f0f7ff;
        }
        .qc-table thead th {
            background: #dbeafe;
            color: #1e40af;
        }

        .audit-table {
            width: 100%;
            border-collapse: separate;
            border-spacing: 0;
            border: 1px solid var(--color-border);
            border-radius: var(--radius-md);
            overflow: hidden;
        }
        .audit-table thead th {
            background: #f7fafc;
            color: var(--color-text-secondary);
            font-size: 11px;
            font-weight: 600;
            letter-spacing: 1px;
            text-transform: uppercase;
            padding: 14px 20px;
            text-align: left;
            border-bottom: 2px solid var(--color-border);
        }
        .audit-table tbody tr {
            transition: background 0.15s ease;
        }
        .audit-table tbody tr:hover {
            background: #f7fafc;
        }
        .audit-table tbody tr:not(:last-child) td {
            border-bottom: 1px solid var(--color-border);
        }
        .audit-table td {
            padding: 14px 16px;
            font-size: 13px;
            vertical-align: top;
            word-break: break-word;
            line-height: 1.7;
        }
        .audit-table td:first-child {
            font-weight: 600;
            color: var(--color-text-primary);
            white-space: nowrap;
            width: 10%;
        }
        .audit-table td:nth-child(2) {
            width: 8%;
            text-align: center;
        }
        .audit-table td:nth-child(3) {
            color: var(--color-text-secondary);
            width: 42%;
        }
        .audit-table td:nth-child(4) {
            color: var(--color-text-secondary);
            width: 40%;
        }

        /* Status badges */
        .status-badge {
            display: inline-flex;
            align-items: center;
            gap: 6px;
            padding: 4px 12px;
            border-radius: 20px;
            font-size: 12px;
            font-weight: 600;
            letter-spacing: 0.5px;
            white-space: nowrap;
        }
        .status-badge::before {
            content: '';
            width: 6px;
            height: 6px;
            border-radius: 50%;
        }
        .status-Pass {
            background: var(--color-pass-bg);
            color: var(--color-pass);
            border: 1px solid var(--color-pass-border);
        }
        .status-Pass::before { background: var(--color-pass); }
        .status-Warning {
            background: var(--color-warning-bg);
            color: var(--color-warning);
            border: 1px solid var(--color-warning-border);
        }
        .status-Warning::before { background: var(--color-warning); }
        .status-Fail {
            background: var(--color-fail-bg);
            color: var(--color-fail);
            border: 1px solid var(--color-fail-border);
        }
        .status-Fail::before { background: var(--color-fail); }

        /* ===== FOOTER ===== */
        .report-footer {
            margin-top: 32px;
            padding-top: 24px;
            border-top: 1px solid var(--color-border);
            display: flex;
            justify-content: space-between;
            align-items: center;
            font-size: 11px;
            color: var(--color-text-muted);
        }
        .footer-logo {
            font-weight: 600;
            letter-spacing: 1px;
        }

        /* ===== RESPONSIVE ===== */
        @media (max-width: 768px) {
            body { padding: 16px; }
            .report-header { padding: 32px 24px; }
            .report-body { padding: 0 24px 32px; }
            .header-title { font-size: 22px; }
            .patient-section { grid-template-columns: 1fr; gap: 16px; }
            .global-summary { flex-wrap: wrap; }
            .case-summary { flex-wrap: wrap; }
            .audit-table td, .audit-table th { padding: 12px 14px; }
            .audit-table td:first-child { white-space: normal; }
        }

        /* ===== PRINT ===== */
        @media print {
            body { background: white; padding: 0; }
            .report-wrapper { max-width: 100%; }
            .report-header { background: #1e3a5f !important; -webkit-print-color-adjust: exact; print-color-adjust: exact; }
            .report-body { box-shadow: none; }
            .audit-table tbody tr:hover { background: transparent; }
            .case-section { page-break-inside: avoid; }
            .no-print { display: none !important; }
        }
    </style>
</head>
<body>
<div class="no-print" style="position: fixed; top: 20px; right: 20px; z-index: 1000;">
    <button onclick="window.print()" style="
        background: linear-gradient(135deg, #1e3a5f 0%, #2b6cb0 100%);
        color: white;
        border: none;
        padding: 12px 24px;
        border-radius: 8px;
        font-size: 14px;
        font-weight: 600;
        cursor: pointer;
        box-shadow: 0 4px 12px rgba(43, 108, 176, 0.3);
        display: flex;
        align-items: center;
        gap: 8px;
        transition: all 0.2s ease;
    " onmouseover="this.style.transform='translateY(-2px)';this.style.boxShadow='0 6px 16px rgba(43, 108, 176,0.4)'" onmouseout="this.style.transform='translateY(0)';this.style.boxShadow='0 4px 12px rgba(43, 108, 176,0.3)'">
        <svg width="16" height="16" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round">
            <path d="M6 9V2h12v7"></path>
            <path d="M6 18H4a2 2 0 0 1-2-2v-5a2 2 0 0 1 2-2h16a2 2 0 0 1 2 2v5a2 2 0 0 1-2 2h-2"></path>
            <rect x="6" y="14" width="12" height="8"></rect>
        </svg>
        导出PDF
    </button>
</div>
<div class="report-wrapper">

    <!-- Header -->
    <header class="report-header">
        <div class="header-label">Medical Compliance Audit</div>
        <h1 class="header-title">病例医学合规批量审核报告</h1>
        <p class="header-subtitle">由 MedAudit Agent 自动生成 · 共 {{ cases|length }} 个病例 · 基于药品规则库的多维度合规审查</p>
    </header>

    <!-- Body -->
    <main class="report-body">

        <!-- Global Summary -->
        <div class="global-summary">
            {% set total_pass = cases|map(attribute='pass_count')|sum %}
            {% set total_warn = cases|map(attribute='warning_count')|sum %}
            {% set total_fail = cases|map(attribute='fail_count')|sum %}
            <div class="stat-chip total"><span class="stat-dot"></span>病例 {{ cases|length }}</div>
            <div class="stat-chip pass"><span class="stat-dot"></span>Pass {{ total_pass }}</div>
            <div class="stat-chip warning"><span class="stat-dot"></span>Warning {{ total_warn }}</div>
            <div class="stat-chip fail"><span class="stat-dot"></span>Fail {{ total_fail }}</div>
        </div>

        {% for case in cases %}
        <section class="case-section">
            <!-- Case Header -->
            <div class="case-header">
                <span class="case-number">{{ loop.index }}</span>
                <span class="case-title">{{ case.filename }}</span>
            </div>

            <!-- Patient Info -->
            <div class="patient-section">
                <div class="info-block">
                    <span class="info-label">患者编号</span>
                    <span class="info-value">{{ case.patient.patient_id }}</span>
                </div>
                <div class="info-block">
                    <span class="info-label">性别 / 年龄</span>
                    <span class="info-value">{{ case.patient.gender }} / {{ case.patient.age }} 岁</span>
                </div>
                <div class="info-block">
                    <span class="info-label">主要诊断</span>
                    <span class="info-value">{{ case.patient.diagnosis }}</span>
                </div>
                <div class="info-block">
                    <span class="info-label">治疗药物</span>
                    <span class="info-value drug">{{ case.patient.treatment_drug }}</span>
                </div>
                <div class="info-block">
                    <span class="info-label">适应症</span>
                    <span class="info-value indication">{{ case.patient.indication }}</span>
                </div>
            </div>

            {% if case.patient.privacy_issues and case.patient.privacy_issues != '无' %}
            <div class="privacy-alert">
                <span class="alert-icon">🔒</span>
                <div class="alert-content">
                    <div class="alert-title">隐私泄露风险检测</div>
                    <div>{{ case.patient.privacy_issues }}</div>
                </div>
            </div>
            {% endif %}

            <!-- Case Summary Stats -->
            <div class="case-summary">
                <div class="stat-chip pass"><span class="stat-dot"></span>Pass {{ case.pass_count }}</div>
                <div class="stat-chip warning"><span class="stat-dot"></span>Warning {{ case.warning_count }}</div>
                <div class="stat-chip fail"><span class="stat-dot"></span>Fail {{ case.fail_count }}</div>
            </div>

            <!-- Audit Results -->
            <h2 class="section-title">详细审核意见</h2>

            {% set qc_aspects = ['文本规范性检查', '通用名规范检查', '医学逻辑审核', '单位规范检查'] %}
            {% set qc_results = case.audit_results|selectattr('aspect', 'in', qc_aspects)|list %}
            {% set other_results = case.audit_results|rejectattr('aspect', 'in', qc_aspects)|list %}

            {% if qc_results %}
            <h3 class="subsection-title">🔍 基础质控检查</h3>
            <table class="audit-table qc-table">
                <thead>
                    <tr>
                        <th style="width: 10%;">检查项目</th>
                        <th style="width: 8%;">状态</th>
                        <th style="width: 42%;">检查详情</th>
                        <th style="width: 40%;">修改建议</th>
                    </tr>
                </thead>
                <tbody>
                    {% for item in qc_results %}
                    <tr>
                        <td>{{ item.aspect }}</td>
                        <td><span class="status-badge status-{{ item.status }}">{{ item.status }}</span></td>
                        <td>{{ item.detail }}</td>
                        <td>{{ item.suggestion }}</td>
                    </tr>
                    {% endfor %}
                </tbody>
            </table>
            {% endif %}

            {% if other_results %}
            <h3 class="subsection-title">📋 合规专项审核</h3>
            <table class="audit-table">
                <thead>
                    <tr>
                        <th style="width: 10%;">审核维度</th>
                        <th style="width: 8%;">状态</th>
                        <th style="width: 42%;">依据详情</th>
                        <th style="width: 40%;">优化建议</th>
                    </tr>
                </thead>
                <tbody>
                    {% for item in other_results %}
                    <tr>
                        <td>{{ item.aspect }}</td>
                        <td><span class="status-badge status-{{ item.status }}">{{ item.status }}</span></td>
                        <td>{{ item.detail }}</td>
                        <td>{{ item.suggestion }}</td>
                    </tr>
                    {% endfor %}
                </tbody>
            </table>
            {% endif %}
        </section>
        {% endfor %}

        <!-- Footer -->
        <footer class="report-footer">
            <span class="footer-logo">MedAudit Agent</span>
            <span>本报告由 AI 自动生成，仅供合规审核参考</span>
        </footer>

    </main>
</div>
</body>
</html>
"""


def generate_combined_report(cases, output_path):
    """利用 Jinja2 渲染所有病例的合并审核报告"""
    template = Template(COMBINED_HTML_TEMPLATE)
    html_out = template.render(cases=cases)

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html_out)
    logger.info(f"合并审核报告已生成: {output_path}")


# ==========================================
# 4. 单个 PPT 处理流程
# ==========================================
# 日志锁，确保并行时每个文件的日志连续输出
import threading
_log_lock = threading.Lock()

def process_single_ppt(pptx_path, output_dir):
    """处理单个 PPT 文件的完整审核流程"""
    filename = os.path.splitext(os.path.basename(pptx_path))[0]
    short_name = filename[:50] + ".." if len(filename) > 50 else filename
    log_lines = []  # 收集本文件的所有日志

    log_lines.append(f"▶ 开始处理: {short_name}")

    # 第一步：提取文本
    log_lines.append("  ├─ [1/4] 解析PPT文本...")
    raw_text = extract_text_from_pptx(pptx_path)

    # 第二步：提取图片并识别
    log_lines.append("  ├─ [2/4] 提取图片...")
    image_info_list, img_log = extract_images_from_pptx(pptx_path, output_dir)
    log_lines.append(img_log)
    image_desc = ""
    if image_info_list:
        if VISION_ENABLED:
            log_lines.append(f"  │  调用视觉模型识别...")
            image_desc = describe_images_with_llm(image_info_list)
            log_lines.append("  │  图片识别完成")
    else:
        log_lines.append("  │  未发现嵌入图片")

    # 第三步：结构化
    log_lines.append("  ├─ [3/4] LLM结构化患者数据...")
    patient_json = transform_raw_text_to_structured(raw_text, image_desc)

    drug = patient_json['treatment_drug']
    indication = patient_json.get('indication', '未识别')
    log_lines.append(f"  │  用药: {drug} | 适应症: {indication}")

    if patient_json.get('privacy_issues') and patient_json['privacy_issues'] != '无':
        log_lines.append(f"  │  ⚠ 隐私风险: {patient_json['privacy_issues'][:60]}...")

    # 第四步：审核
    log_lines.append("  └─ [4/4] 合规审核...")
    rules = load_rules(drug, indication)
    audit_results = audit_patient_data(patient_json, rules)

    pass_n = sum(1 for r in audit_results if r['status'] == 'Pass')
    warn_n = sum(1 for r in audit_results if r['status'] == 'Warning')
    fail_n = sum(1 for r in audit_results if r['status'] == 'Fail')
    log_lines.append(f"     结果: ✓{pass_n} Pass | ⚠{warn_n} Warning | ✗{fail_n} Fail")
    log_lines.append("")

    # 使用锁一次性输出所有日志，避免交错
    with _log_lock:
        for line in log_lines:
            logger.info(line)

    return {
        "file": filename,
        "drug": drug,
        "indication": indication,
        "images": len(image_info_list),
        "audit_count": len(audit_results),
        "pass_count": pass_n,
        "warning_count": warn_n,
        "fail_count": fail_n,
        "patient_data": patient_json,
        "audit_results": audit_results,
    }


# ==========================================
# 主运行流 - 批量处理 + 并行加速
# ==========================================
if __name__ == "__main__":
    input_dir = "data/input_cases"
    output_dir = "data/output_reports"
    os.makedirs(output_dir, exist_ok=True)

    start_time = time.time()

    logger.info("=" * 60)
    logger.info("  MedAudit Agent - 医学病例合规自动审核系统")
    logger.info("=" * 60)
    logger.info(f"  输入目录: {input_dir}")
    logger.info(f"  输出目录: {output_dir}")
    logger.info(f"  模型: {MODEL_NAME} | 视觉: {'已启用' if VISION_ENABLED else '未启用'}")
    logger.info("-" * 60)

    # 扫描所有 PPT 文件
    all_ppt_files = sorted(glob.glob(os.path.join(input_dir, "*.pptx")))

    if not all_ppt_files:
        logger.warning(f"⚠ 未在 {input_dir} 目录下找到任何 .pptx 文件")
    else:
        logger.info(f"  发现 {len(all_ppt_files)} 个PPT文件，开始批量审核...")
        logger.info("")

        # 并行处理（最多3个线程，避免API限流）
        results = []
        max_workers = min(3, len(all_ppt_files))

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_file = {
                executor.submit(process_single_ppt, ppt_file, output_dir): ppt_file
                for ppt_file in all_ppt_files
            }

            for future in as_completed(future_to_file):
                ppt_file = future_to_file[future]
                try:
                    result = future.result()
                    results.append(result)
                except Exception as e:
                    filename = os.path.basename(ppt_file)
                    logger.error(f"✗ 处理失败: {filename}")
                    logger.error(f"  错误: {e}")
                    results.append({"file": filename, "error": str(e)})

        # 按原始文件顺序排序结果
        file_order = {os.path.basename(f): i for i, f in enumerate(all_ppt_files)}
        results.sort(key=lambda r: file_order.get(r["file"], 999))

        # 生成合并报告
        success_results = [r for r in results if "error" not in r]
        if success_results:
            combined_cases = [
                {
                    "filename": r["file"],
                    "patient": r["patient_data"],
                    "audit_results": r["audit_results"],
                    "pass_count": r["pass_count"],
                    "warning_count": r["warning_count"],
                    "fail_count": r["fail_count"],
                }
                for r in success_results
            ]
            combined_path = os.path.join(output_dir, "audit_report_combined.html")
            generate_combined_report(combined_cases, combined_path)

        # 汇总统计
        elapsed_time = time.time() - start_time
        total_pass = sum(r.get("pass_count", 0) for r in success_results)
        total_warn = sum(r.get("warning_count", 0) for r in success_results)
        total_fail = sum(r.get("fail_count", 0) for r in success_results)
        fail_files = [r for r in results if "error" in r]

        logger.info("")
        logger.info("=" * 60)
        logger.info("  处理完成")
        logger.info("=" * 60)
        logger.info(f"  耗时: {elapsed_time:.1f}秒 | 成功: {len(success_results)}/{len(results)}")
        logger.info(f"  审核统计: Pass={total_pass} | Warning={total_warn} | Fail={total_fail}")
        logger.info("-" * 60)

        # 详细结果表格
        logger.info("")
        logger.info("  详细结果:")
        logger.info("  " + "-" * 90)
        logger.info(f"  {'文件名':<40} {'药品':<12} {'适应症':<8} {'Pass':>4} {'Warn':>4} {'Fail':>4}")
        logger.info("  " + "-" * 90)

        for r in results:
            if "error" in r:
                short_name = r['file'][:38] + ".." if len(r['file']) > 40 else r['file']
                logger.info(f"  {short_name:<40} {'❌ 失败':<12}")
                logger.info(f"  {'':>40} 原因: {r['error'][:50]}")
            else:
                short_name = r['file'][:38] + ".." if len(r['file']) > 40 else r['file']
                logger.info(f"  {short_name:<40} {r['drug']:<12} {r['indication']:<8} {r['pass_count']:>4} {r['warning_count']:>4} {r['fail_count']:>4}")

        logger.info("  " + "-" * 90)
        logger.info("")

        if fail_files:
            logger.info(f"  ⚠ {len(fail_files)} 个文件处理失败，请检查日志了解详情")

        logger.info(f"  📄 合并报告: {output_dir}/audit_report_combined.html")
        logger.info(f"  📋 详细日志: {LOG_DIR}/medaudit.log")
        logger.info("=" * 60)
