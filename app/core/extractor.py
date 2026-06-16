"""
Agent 1: PPT 文本 + 图片提取
从 .pptx 文件中提取文本内容和嵌入图片
"""

import os
import base64
import logging

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.config import settings
from app.services.llm_client import get_sync_client

logger = logging.getLogger("MedAudit")


def extract_text_from_pptx(pptx_path: str) -> str:
    """提取PPT内所有文本框和表格数据"""
    logger.debug(f"开始提取PPT文本: {pptx_path}")
    prs = pptx.Presentation(pptx_path)
    extracted_texts = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_text = f"--- Slide {slide_idx + 1} --- \n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_text += paragraph.text.strip() + "\n"
            if shape.has_table:
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_data:
                        slide_text += " [表格数据]: " + " | ".join(row_data) + "\n"
        extracted_texts.append(slide_text)

    logger.debug(f"提取完成，共 {len(prs.slides)} 页幻灯片")
    return "\n".join(extracted_texts)


def extract_images_from_pptx(pptx_path: str, output_dir: str) -> tuple[list, str]:
    """
    提取PPT中所有嵌入的图片，保存到指定目录

    Returns:
        (image_info_list, log_line)
        image_info_list: [(img_path, slide_idx), ...]
    """
    logger.debug(f"开始提取PPT图片: {pptx_path}")
    prs = pptx.Presentation(pptx_path)
    filename = os.path.splitext(os.path.basename(pptx_path))[0]
    img_dir = os.path.join(output_dir, "images", filename)
    os.makedirs(img_dir, exist_ok=True)

    image_info_list = []
    skipped_count = 0

    for slide_idx, slide in enumerate(prs.slides):
        img_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_count += 1
                    image = shape.image
                    ext = image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    # 跳过不支持的格式
                    if ext.lower() not in ['bmp', 'gif', 'jpg', 'png', 'tiff', 'wmf']:
                        skipped_count += 1
                        logger.debug(f"  跳过不支持的图片格式: {ext}")
                        continue
                    img_filename = f"slide_{slide_idx + 1}_img_{img_count}.{ext}"
                    img_path = os.path.join(img_dir, img_filename)

                    with open(img_path, "wb") as f:
                        f.write(image.blob)

                    image_info_list.append((img_path, slide_idx + 1))
                    logger.debug(f"  保存图片: {img_filename}")
                except ValueError as e:
                    skipped_count += 1
                    logger.debug(f"  跳过图片: {e}")
                    continue

    if skipped_count > 0:
        log_line = f"  │  提取图片 {len(image_info_list)} 张（跳过 {skipped_count} 张不支持的格式）"
    else:
        log_line = f"  │  提取图片 {len(image_info_list)} 张"
    return image_info_list, log_line


def image_to_base64(image_path: str) -> str:
    """将图片文件转为 base64 编码"""
    with open(image_path, "rb") as f:
        return base64.b64encode(f.read()).decode("utf-8")


def describe_images_with_llm(image_info_list: list) -> str:
    """利用多模态 LLM 识别图片内容，返回描述文本"""
    if not image_info_list or not settings.vision_enabled:
        if image_info_list and not settings.vision_enabled:
            logger.warning("    视觉模型未配置，跳过图片识别。设置 OPENAI_MULTIMODAL_MODEL 环境变量以启用。")
        return ""

    client = get_sync_client()
    descriptions = []

    for img_path, slide_idx in image_info_list:
        try:
            img_b64 = image_to_base64(img_path)
            ext = os.path.splitext(img_path)[1].lower()
            mime_map = {
                ".png": "image/png", ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg", ".gif": "image/gif", ".webp": "image/webp"
            }
            mime_type = mime_map.get(ext, "image/png")

            logger.debug(f"  调用视觉模型识别: Slide {slide_idx} 图片")
            response = client.chat.completions.create(
                model=settings.openai_multimodal_model,
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "请仔细识别这张医学病例PPT中的图片内容。如果是化验单/检查报告，请提取关键数值（如EOS、FEV1等）；如果是CT/MRI影像，请描述所见；如果包含患者个人信息（姓名、住院号、医生信息），请特别标注。用中文简洁描述。"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{img_b64}"
                            }
                        }
                    ]
                }],
                temperature=0.1
            )
            desc = response.choices[0].message.content.strip()
            descriptions.append(f"[Slide {slide_idx} 图片内容]: {desc}")
            logger.debug(f"  Slide {slide_idx} 图片识别完成: {desc[:50]}...")
        except Exception as e:
            logger.error(f"  Slide {slide_idx} 图片识别失败: {e}")
            descriptions.append(f"[Slide {slide_idx} 图片识别失败]: {str(e)}")

    return "\n".join(descriptions)
